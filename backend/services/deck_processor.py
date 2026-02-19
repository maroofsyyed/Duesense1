"""
Deck Processor for DueSense
Extracts text from PDF/PPTX pitch decks and structures with LLM.
"""
import io
import logging
from pypdf import PdfReader
from pptx import Presentation
from services.llm_provider import llm

logger = logging.getLogger(__name__)


async def extract_deck(file_path: str, file_ext: str) -> dict:
    """
    Extract structured data from a pitch deck file.
    
    Args:
        file_path: Path to the uploaded file
        file_ext: File extension (pdf, pptx, ppt)
        
    Returns:
        Structured data dict with company info, founders, etc.
        
    Raises:
        ValueError: If file cannot be read or has no extractable text
        RuntimeError: If LLM processing fails
    """
    logger.info(f"📄 Starting deck extraction: {file_path} (type: {file_ext})")
    
    # Read file content
    try:
        with open(file_path, "rb") as f:
            content = f.read()
        logger.info(f"✓ File read: {len(content):,} bytes")
    except Exception as e:
        logger.error(f"❌ Failed to read file {file_path}: {e}")
        raise ValueError(f"Could not read file: {e}")

    # Extract text based on file type
    try:
        if file_ext == "pdf":
            text = _extract_pdf(content)
        else:
            text = _extract_pptx(content)
        logger.info(f"✓ Text extracted: {len(text):,} chars")
    except Exception as e:
        logger.error(f"❌ Text extraction failed: {type(e).__name__}: {e}")
        raise ValueError(f"Could not extract text from {file_ext.upper()} file: {e}")

    # Validate extracted text
    if not text or len(text.strip()) < 50:
        logger.error(f"❌ Insufficient text extracted: {len(text) if text else 0} chars")
        raise ValueError(
            f"Could not extract meaningful text from the file. "
            f"The file may be image-based, corrupted, or password-protected. "
            f"Extracted {len(text) if text else 0} characters."
        )

    # Structure with LLM
    try:
        logger.info(f"🤖 Sending {len(text[:12000])} chars to LLM for structuring...")
        structured = await _structure_with_llm(text)

        # Validate extraction — company name must exist
        company_name = structured.get("company", {}).get("name")
        if not company_name or company_name in ("", "Unknown", "Unknown Company", "null"):
            logger.warning("⚠️ Company name not extracted — retrying with simplified prompt...")
            structured = await _structure_with_llm(text)
            company_name = structured.get("company", {}).get("name")

        if not company_name or company_name in ("", "Unknown", "Unknown Company", "null"):
            logger.error("❌ Company name still not extracted after retry")
            # Still return what we have — partial data is better than nothing
        else:
            logger.info(f"✓ Extracted company: {company_name}")

        logger.info(f"✓ LLM structuring complete — keys: {list(structured.keys())}")
        return structured
    except Exception as e:
        logger.error(f"❌ LLM structuring failed: {type(e).__name__}: {e}")
        raise RuntimeError(f"AI analysis failed: {e}")


def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF content."""
    try:
        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)
        parts = []
        
        logger.info(f"  Processing {len(reader.pages)} PDF pages...")
        
        for page_num, page in enumerate(reader.pages, 1):
            try:
                t = page.extract_text()
                if t:
                    parts.append(t)
            except Exception as page_err:
                logger.warning(f"  ⚠️ Failed to extract page {page_num}: {page_err}")
                continue
        
        text = "\n\n".join(parts)
        logger.info(f"  Extracted text from {len(parts)}/{len(reader.pages)} pages")
        return text
        
    except Exception as e:
        logger.error(f"  ❌ PDF extraction error: {e}")
        raise ValueError(f"PDF extraction failed: {e}")


def _extract_pptx(content: bytes) -> str:
    """Extract text from PPTX content."""
    try:
        pptx_file = io.BytesIO(content)
        prs = Presentation(pptx_file)
        parts = []
        
        logger.info(f"  Processing {len(prs.slides)} slides...")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
            except Exception as shape_err:
                logger.warning(f"  ⚠️ Error on slide {slide_num}: {shape_err}")
            parts.append("\n".join(slide_text))
        
        text = "\n\n".join(parts)
        logger.info(f"  Extracted text from {len(parts)} slides")
        return text
        
    except Exception as e:
        logger.error(f"  ❌ PPTX extraction error: {e}")
        raise ValueError(f"PPTX extraction failed: {e}")


async def _structure_with_llm(text: str) -> dict:
    prompt = f"""Extract structured data from this startup pitch deck.

OUTPUT FORMAT: Return ONLY a valid JSON object. No markdown, no code blocks, no explanations.

CRITICAL JSON RULES:
- Use null (not "null", not "not_mentioned", not "N/A") for missing values
- Use true/false (not "true"/"false" strings) for booleans
- All strings must be in double quotes
- Arrays can be empty [] if no data found

EXTRACTION RULES:
- ONLY extract information EXPLICITLY stated in the text
- Use null for any missing fields - NEVER guess or hallucinate
- Quote exact numbers from the deck - no rounding

JSON STRUCTURE TO FILL:
{{
  "company": {{
    "name": "string",
    "tagline": "string or null",
    "founded": "string year or null",
    "hq_location": "string or null",
    "website": "string or null",
    "stage": "pre-seed|seed|series-a|series-b|series-c+|null",
    "industry": "string"
  }},
  "founders": [
    {{
      "name": "string",
      "role": "string",
      "linkedin": "string or null",
      "github": "string or null",
      "previous_companies": [],
      "years_in_industry": null,
      "education": "string or null"
    }}
  ],
  "problem": {{
    "statement": "string",
    "market_pain": "string or null",
    "current_solutions": []
  }},
  "solution": {{
    "product_description": "string",
    "key_features": [],
    "technology_stack": [],
    "ai_usage": {{
      "is_ai_core": false,
      "ai_description": null,
      "proprietary_data": false,
      "model_architecture": null
    }}
  }},
  "market": {{
    "tam": "string with $ or null",
    "sam": "string with $ or null",
    "som": "string with $ or null",
    "growth_rate": "string or null",
    "target_customers": "string or null"
  }},
  "traction": {{
    "revenue": "string with $ or null",
    "mrr": "string or null",
    "customers": null,
    "growth_rate": "string or null",
    "key_metrics": {{}}
  }},
  "business_model": {{
    "type": "SaaS|Marketplace|Enterprise|Other|null",
    "pricing": "string or null",
    "unit_economics": "string or null"
  }},
  "funding": {{
    "seeking": "string with $ or null",
    "previous_rounds": [],
    "total_raised": "string or null",
    "valuation": "string or null"
  }},
  "competitive_advantages": [],
  "risks": []
}}

PITCH DECK TEXT:
{text[:12000]}

OUTPUT (valid JSON only):"""

    system = """You are a precise JSON data extraction assistant. 
IMPORTANT: Output ONLY valid JSON. No explanations, no markdown.
Use null for missing values. Use true/false for booleans.
Never use unquoted words like not_mentioned or N/A."""

    result = await llm.generate_json(prompt, system)
    return result
